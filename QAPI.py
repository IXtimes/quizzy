import json
import fitz, docx2txt
from pptx import Presentation
from cryptography.fernet import Fernet
from tkinter import filedialog
import Levenshtein as lev
from openai import OpenAI
import requests
from prompts import *
from keys import CONTEXT_CRYPTO_KEY
from tkinter import messagebox
import random
import threading
import time

CONTEXT_CUTOFF = 350
THREAD_TIMEOUT_TIME = 20
MAX_THREADS = 5
THREAD_SIZE = 1

def call_gpt(domain, context, api_key, model, system, request, image_embed):
    return_response = []
    
    # create a background function to make the api call
    def internal_gpt_call():
        # initalize client and domain/context segments
        client = OpenAI(api_key=api_key)
        domain_c = "Domain: " + domain + "\n"
        context_c = "Context: " + get_random_context_segment(context, CONTEXT_CUTOFF) + "\n"
        model_c = determine_model_str_from_index(model)
        
        try:
            response = client.chat.completions.create (
            model = model_c,
            messages = [
                    {"role": "system", "content": system},
                    {
                        "role":"user",
                        "content": [
                            {"type": "text", "text": domain_c + context_c + request},
                            {"type": "image_url", "image_url": {
                                "url": image_embed
                            }} 
                        ] if image_embed != "" else [
                            {"type": "text", "text": domain_c + context_c + request}
                        ]
                    }
                ]
            )
            return_response.append(response.choices[0].message.content)
        except Exception as e:
            return_response.append(f"Error: {str(e)}")
        
    # push the response to a thread and wait for it using a timeout length
    thread = threading.Thread(target=internal_gpt_call)
    thread.start()
    thread.join(THREAD_TIMEOUT_TIME)
    
    # check if the thread timed out
    if thread.is_alive():
        raise TimeoutError("GPT call timed out!")

    return return_response[0]

def determine_model_str_from_index(model):
    return "gpt-4o-mini" if model == 0 else ("gpt-4o" if model == 1 else "o1-mini")

def is_number(s):
    try:
        float(s)
        return True
    except ValueError:
        return False

def all_at_once_numaric_comparision_grading(answers, correct_answers):
    # get a list of similarities and copy the collections
    similarities = []
    ans = answers.copy()
    cor_ans = correct_answers.copy()
    
    # iterate through the numbers in the ans list
    for answer in ans:
        # check if this answer is numaric
        if not is_number(answer):
            similarities.append(0.0)
            
            continue
        
        most_similar = ""
        best_similarity = 0.0
        # iterate through the cor_ans list
        for correct in cor_ans:
            # get the percent distance from this answer
            distance = abs(float(answer) - float(correct)) / float(correct)
            # get a ratio of value with a 5% confidence
            similarity = 1 - min(distance / 0.05, 1)
            
            # compare with the score already logged for this answer
            if similarity > best_similarity:
                # save as best, remember with what string
                best_similarity = similarity
                most_similar = correct
                
        # check if a record was found to be similar
        if best_similarity == 0.0:
            # append a score of 0
            similarities.append(0.0)
            
            # continue
            continue
                
        # record this similarity and remove the most similar string from answer collection
        similarities.append(best_similarity)
        cor_ans.remove(most_similar)
        
    # return results
    return similarities

def all_at_once_answer_comparision_grading(answers, correct_answers):
    # get a list of similarities and copy the lowercase conversion collections
    similarities = []
    ans = [ans.lower() for ans in answers]
    cor_ans = [ans.lower() for ans in correct_answers]
    
    # iterate through the strings in the ans list
    for answer in ans:
        most_similar = ""
        best_similarity = 0.0
        best_diff = 0.0
        
        # iterate through the cor_ans list
        for correct in cor_ans:
            # get the Levenshtein similarity score between these 2 strinsg
            distance = lev.distance(answer, correct)
            max_length = max(len(answer), len(correct))
            difference = (distance / max_length)
            # get a ratio of value with 25% confidence
            similarity = 1 - min(difference/.25, 1)
            
            # compare with the score already logged for this answer
            if similarity > best_similarity:
                # save as best, remember with what string
                best_similarity = similarity
                most_similar = correct
                best_diff = difference
                
        # check if a record was found to be similar
        if best_similarity == 0.0:
            # append a score of 0
            similarities.append(0.0)
            
            # continue
            continue
                
        # record this similarity and remove the most similar string from answer collection and this answer
        similarities.append(best_similarity)
        cor_ans.remove(most_similar)
        
    # return results
    return similarities

def answer_comparision_grading(answer, correct_answer):
    # get the Levenshtein similarity score between these 2 strinsg
    distance = lev.distance(answer, correct_answer)
    max_length = max(len(answer), len(correct_answer))
    difference = (distance / max_length)
    # get a ratio of value with 25% confidence
    similarity = 1 - min(difference/.25, 1)
    
    # return this similarity
    return similarity

def validate_question_object(object):
    # ensure we have a question text and q_index, otherwise fail immediately
    if not 'Question' in object or not 'Q' in object:
        print("@FAIL (???): Object too corrupted to be a valid question object. Terminating...")
        messagebox.showerror("Failure @?", "Unable to parse question header. File might be corrupted!")
        return False
    
    # for debugging, grab the question's index
    q_index = object['Q']
    
    # determine the question's type
    try:
        if not 'Type' in object:
            print(f'@WARN ({q_index}): Question missing its type, defaulting to MC!')
            object['Type'] = 'MC'
        match object['Type']:
            case "MC":
                # count the incorrect and correct answers and ensure they enumerate correctly
                num_cor = 0
                max_cor = 0
                num_incor = 0
                max_incor = 0
                
                # iterate over the keys
                for key, ___ in object.items():
                    if not key.startswith(("C", "A")):
                        continue
                    k_prefix = key[0]
                    k_val = int(key[1:])
                    
                    # check for a correct key
                    if k_prefix == "C":
                        # count this key and log if this key is the largest found so far
                        num_cor += 1
                        max_cor = max(max_cor, k_val)
                    # check for a incorrect key
                    if k_prefix == "A":
                        # count this key and log if this key is the largest found so far
                        num_incor += 1
                        max_incor = max(max_incor, k_val)
                        
                # check if we are missing correct answers or incorrect answers entirely
                if num_cor == 0:
                    print(f'@FAIL ({q_index}): Question missing correct answers. Terminating...')
                    messagebox.showerror(f"Failure @{q_index}", "Question missing correct answers, unable to import!")
                    return False  
                if num_incor == 0:
                    print(f'@FAIL ({q_index}): Question missing incorrect answers. Terminating...')
                    messagebox.showerror(f"Failure @{q_index}", "Question missing incorrect answers, unable to import!")
                    return False  
                    
                # check if the counts match
                if num_cor != max_cor:
                    print(f'@FAIL ({q_index}): Question correct answer count does not match maximal index. Terminating...')
                    messagebox.showerror(f"Failure @{q_index}", "Correct answers are incorrectly indexed, unable to import!")
                    return False
                if num_incor != max_incor:
                    print(f'@FAIL ({q_index}): Question incorrect answer count does not match maximal index. Terminating...')
                    messagebox.showerror(f"Failure @{q_index}", "Incorrect answers are incorrectly indexed, unable to import!")
                    return False
                
                # check for a forced state, default to "Either" (0)
                if not "Forced" in object:
                    print(f'@WARN ({q_index}): Question lacking "Forced" flag, defaulting to Either (0)!')
                    object["Forced"] = "0"
                
                # check for illegal keys
                for key, ___ in object.items():
                    if key == "Type":
                        continue
                    if key.startswith(("T", "D", "Guidelines", "Format", "Language")):
                        print(f'@FAIL ({q_index}): MC question contains illegal key: "{key}" Terminating...')
                        messagebox.showerror(f"Failure @{q_index}", f"Multiple choice question contains illegal key {key}, unable to import!")
                        return False
            case "TD":
                # count the terms and definitions and ensure they enumerate correctly
                num_terms = 0
                max_terms = 0
                num_defines = 0
                max_defines = 0
                
                # iterate over the keys
                for key, ___ in object.items():
                    if not key.startswith(("T", "D")) or key.startswith("Ty"):
                        continue
                    k_prefix = key[0]
                    k_val = int(key[1:])
                    
                    # check for a term key
                    if k_prefix == "T":
                        # count this key and log if this key is the largest found so far
                        num_terms += 1
                        max_terms = max(max_terms, k_val)
                    # check for a definition key
                    if k_prefix == "D":
                        # count this key and log if this key is the largest found so far
                        num_defines += 1
                        max_defines = max(max_defines, k_val)
                        
                # check if we are missing terms or definitions entirely
                if num_terms == 0:
                    print(f'@FAIL ({q_index}): Question missing terms. Terminating...')
                    messagebox.showerror(f"Failure @{q_index}", "Question is missing term keys, unable to import!")
                    return False  
                if num_defines == 0:
                    print(f'@FAIL ({q_index}): Question missing definitions. Terminating...')
                    messagebox.showerror(f"Failure @{q_index}", "Question is missing definition keys, unable to import!")
                    return False  
                        
                # check if the counts match
                if num_terms != max_terms:
                    print(f'@FAIL ({q_index}): Question term count does not match maximal index. Terminating...')
                    messagebox.showerror(f"Failure @{q_index}", "Terms are incorrectly indexed, unable to import!")
                    return False
                if num_defines != max_defines:
                    print(f'@FAIL ({q_index}): Question definition count does not match maximal index. Terminating...')
                    messagebox.showerror(f"Failure @{q_index}", "Definitions are incorrectly indexed, unable to import!")
                    return False
                
                # check for a forced state, default to "Either" (0)
                if not "Forced" in object:
                    print(f'@WARN ({q_index}): Question lacking "Forced" flag, defaulting to Either (0)!')
                    object["Forced"] = "0"
                
                # check for illegal keys
                for key, ___ in object.items():
                    if key == "Type":
                        continue
                    if key.startswith(("C", "A", "Guidelines", "Format", "Language")):
                        print(f'@FAIL ({q_index}): TD question contains illegal key: "{key}" Terminating...')
                        messagebox.showerror(f"Failure @{q_index}", f"Term-definition question contains illegal key {key}, unable to import!")
                        return False
            case "Ess":
                # fail if we lack guidelines
                if not "Guidelines" in object:
                    print(f'@FAIL ({q_index}): Essay question missing guidelines. Terminating...')
                    messagebox.showerror(f"Failure @{q_index}", "Question is missing grading guidelines, unable to import!")
                    return False
                
                # check for format specification, and default to "0" if missing
                if not "Format" in object:
                    print(f'@WARN ({q_index}): Essay question missing format, defaulting to Explaination (0)!')
                    object["Format"] = "0"
                    
                # check for difficulty specification, and default to "1" if missing
                if not "Difficulty" in object:
                    print(f'@WARN ({q_index}): Essay question missing difficulty, defaulting to Know (1)!')
                    object["Difficulty"] = "1"
                
                # fail if format is "Code" and the language is not specified
                if not "Language" in object or (object["Format"] == "1" and object["Language"] == "N/A"):
                    print(f'@FAIL ({q_index}): Essay question missing language when "Code" was specified. Terminating...')
                    messagebox.showerror(f"Failure @{q_index}", "Question is missing a language specification when 'Code' was specified, unable to import!")
                    return False
                
                # check for illegal keys
                for key, ___ in object.items():
                    if key == "Type":
                        continue
                    if key.startswith(("T", "D", "A", "C")) and key != "Difficulty":
                        print(f'@FAIL ({q_index}): Essay question contains illegal key: "{key}" Terminating...')
                        messagebox.showerror(f"Failure @{q_index}", f"Essay question contains illegal key {key}, unable to import!")
                        return False
            case _:
                # fail for invalid type
                print(f'@FAIL ({q_index}): Question failed validation due to corrupted type: "{object["Type"]}". Terminating...')
                messagebox.showerror(f"Failure @{q_index}", f"{object["Type"]} is not a valid question type, unable to import!")
                return False
                    
    except Exception as e:
        print(f'@FAIL ({q_index}): Question failed validation due to exception: "{e}". Terminating...')
        messagebox.showerror(f"Failure @{q_index}", f"Importing this question generated the following exception: {e}, unable to import!")
        return False
        
    # check for an explaination, default to "No explaination was provided :("
    if not "Explaination" in object:
        print(f'@WARN ({q_index}): Question lacks explaination field, default provided!')
        object["Explaination"] = "No explaination was provided :("
        
    # if we make it here, then this question must be valid!
    return True                

def get_random_context_segment(context, cutoff):
    # split the context segment into words and randomly choose a segment of those words based on the cuttoff
    context_words = context.split(" ")
    start_point = random.choice(range(0, max(len(context_words) - cutoff, 1)))
    return " ".join(context_words[start_point:min(start_point+cutoff, len(context_words))]) + "\n"

def remove_backticked_imbeds(s):
    # split the input based on backticked imbeds
    imbeds = s.split('```')
    result = ""
    
    # iterate through the imbeds
    for i, chunk, in enumerate(imbeds):
        if i % 2 == 0:
            # not an embed, skip
            result += chunk
        else:
            # readd if not an image embed
            if not chunk.startswith("img:"):
                result += "```" + chunk + "```"
    
    return result

def get_image_embed_links_if_any(question):
    # Split the question content into chunks to search for the first img embed
    imbeds = question.split('```')
    
    for i, chunk in enumerate(imbeds):
        if i % 2 != 0:
            # determine what imbed we are using
            header = chunk.split(':')[0]
            if header == "img":
                # we are imbedding an image, remove header from string
                link = chunk[4:]
                
                # fetch the image using a request
                response = requests.get(link)
                
                # check if successful
                if response.status_code == 200:
                    # return link
                    return link
                else:
                    # return blank
                    return ""
                
    # return blank if no link is found
    return ""

def grade_multiple_choice_question(question, answers, frqComp):
    # collect the correct answers of the question
    correct_answers = []
    for key, val in question.items():
        if key.startswith("C"):
            correct_answers += [val]
            
    # determine if we are doing numaric comparision grading
    is_numaric = all([is_number(ans) for ans in correct_answers])
                
    # determine if this question is prosed as a FRQ
    if frqComp:
        max_weighted_score = len(correct_answers)
        score = 0
        
        # compare each answer choice against all possible correct answers and take the one with the best score (answers are written and thus may partially match)
        if is_numaric:
            # get the best matching answer-score pair
            results = all_at_once_numaric_comparision_grading(answers, correct_answers)
            
            # add the results sum to the total score
            score += sum(results)
        else:
            # get the best matching answer-score pair
            results = all_at_once_answer_comparision_grading(answers, correct_answers)
            
            # add the results sum to the total score
            score += sum(results)
            
        # scale the score out of 1
        score = score / max_weighted_score
    else:
        # answers must absolutely match, so we can simply just do a symmetric difference of the 2 sets
        missing_matches = [x for x in list(set(answers).symmetric_difference(set(correct_answers))) if x != "NanX"]
        print(answers)
        print(correct_answers)
        print(missing_matches)
        
        # the score is the inverse of the ratio of missing matches over total correct answers
        score = 1 - (min(len(missing_matches), len(correct_answers)) / len(correct_answers))
        
    return int(score * 10)

def grade_matching_question(question, answers, frqComp):
    # collect the terms in the order that they appear from the question
    matchings = sum(1 for key in question if key.startswith("D"))
    terms = []
    for i in range(matchings):
        terms += [question["T" + str(i + 1)]]
        
    # determine if this question is prosed as a FRQ
    max_weighted_score = len(terms)
    score = 0
    if frqComp:
        # iterate over the answers and compare them to their respective term, where score can be partial if the provided answer is close
        score = sum([answer_comparision_grading(ans, term) for ans, term in zip(answers, terms)]) / max_weighted_score
    else:
        # iterate over the answers and just check if they are equal, no partial credit
        score = sum([1 if ans == term else 0 for ans, term in zip(answers, terms)]) / max_weighted_score
        
    return int(score * 10)

def grade_question(question, Q_index, answers, frqComp, domain, context, api_key, model):
    # grade this question differently depending on its type
    match question['Type']:
        case "MC":
           score = grade_multiple_choice_question(question, answers, frqComp)
        case "TD":
            score = grade_matching_question(question, answers, frqComp)
        case "Ess":
            # score is completely determined by prompt through assessing the answer that the user typed in and how it adheres to the question guidelines
            score = grade_essay_question(question, answers[0], domain, context, api_key, model)
            
    # scale the score to 10 points and return the result
    return int(score * 10)
            
def grade_essay_question(question, answer, domain, context, api_key, model):
    # First retrieve an image embed if it exists
    img_link = get_image_embed_links_if_any(question["Question"])
    
    # determine how the student response will be graded 
    match question['Format']:
        case "0": # Explaination
            format = "Graded as an explaination; "
        case "1": # Code
            format = "Graded as code written in " + question['Language'] + "; "
        case "2": # Proof
            format = "Graded as a proof; "
            
    # determine the difficulty of the question
    match question['Difficulty']:
        case "0": # Easy
            request_difficulty = REQUEST_EXAMPLE_GRADE_ESSAY_GUIDELINES_EASY
        case "1": # Medium
            request_difficulty = REQUEST_EXAMPLE_GRADE_ESSAY_GUIDELINES_MEDIUM
        case "2": # Hard
            request_difficulty = REQUEST_EXAMPLE_GRADE_ESSAY_GUIDELINES_HARD
            
    # Get the response from the GPT
    try:
        system_content = SYSTEM_INTRO + request_difficulty + SYSTEM_CONCLUSION + (IMG_SPECIFICATION if img_link != "" else "")
        response_content = format + "Question: "+ remove_backticked_imbeds(question["Question"]) + "Guidelines: " + question["Guidelines"] + "\n" +  REQUEST_GRADE_ESSAY_GUIDELINES + "\nInput Response: " + answer
        response = call_gpt(domain, context, api_key, model, system_content, response_content, img_link)
        
        # split the response into the grade and explaination
        result = response.strip().replace('\\\\', '\\')
        grade = int(result.split("~")[0])
        explaination = result.split("~")[1]
        
        # override the question's explaination with the one the AI cooked up
        question["Explaination"] = explaination
    except TimeoutError:
        # fail to grade
        messagebox.showerror("GPT Call Fail!", "A call to ChatGPT to grade an essay question timed out! So it will unfortunately not be graded.")
        grade = 10
    
    # return the score
    return int(grade)

def generate_permutation_from_question(question, domain, context, api_key, model, t_results = None):
    # create a new question object
    new_question = {}
    
    # seed the question text with generated permutations (permutation questions dont have an index as they are not saved)
    new_question['Q'] = str(-1)
    new_question['Forced'] = str(1)
    new_question['Question'] = seed_permutations_in_question(question["Question"], domain, context, api_key, model)
    new_question['Type'] = question['Type']

    # from the new question header text, generate new answer choices that emulate the same pattern as the input
    match new_question['Type']:
        case "MC":
            fill_multiple_choice_options(new_question, domain, context, api_key, model)
        case 'TD':
            fill_matching_options(new_question, False, domain, context, api_key, model)
        case 'Ess':
            fill_essay_guidelines(new_question, domain, context, api_key, model)

    # finally, generate the explaination for the question.
    generate_explaination_for_question(new_question, domain, context, api_key, model)

    # push results to threads, otherwise return
    if t_results is None:
        return new_question
    else:
        t_results.append(new_question)

def unscramble_matching(question, domain, context, api_key, model):
    # First retrieve an image embed if it exists
    img_link = get_image_embed_links_if_any(question["Question"])
    
    # get the list of terms and definitions for appropriate matching
    pair_count = sum(1 for key in question if key.startswith("D"))
    terms = []
    definitions = []
    for i in range(pair_count):
        terms += [question["T" + str(i + 1)]]
        definitions += [question["D" + str(i + 1)]]
    
    # Get the response from the GPT
    try:
        system_content = SYSTEM_INTRO + REQUEST_EXAMPLE_UNSCRAMBLE + SYSTEM_CONCLUSION + (IMG_SPECIFICATION if img_link != "" else "")
        response_content = REQUEST_UNSCRAMBLE + "Question: " + remove_backticked_imbeds(question["Question"]) + "Terms: " + "\n".join([f"{i + 1}. {term}" for i, term in enumerate(terms)]) + "\nDefinitions: " + "\n".join([f"{i + 1}. {defin}" for i, defin in enumerate(definitions)])
        response = call_gpt(domain, context, api_key, model, system_content, response_content, img_link)
        
        # the result is a list of definition indexes in the order of how they are supposed to match with the term indexes
        try:
            # using the list, simply iterate in the order of the definitions and update them according to our array of numbers
            matchings = [int(num) for num in response.split(",")]
            for i in range(pair_count):
                question['D' + str(i + 1)] = definitions[matchings[i] - 1]
            
        except Exception as e:
            # report if this generates an exception when parsing
            messagebox.showerror("GPT Error!", "ChatGPT was unable to unscrable these terms and definitions.")
            print("GPT failed the matching: " + str(e))
            
            return
    except TimeoutError:
        # fail to unscramble
        messagebox.showerror("GPT Call Fail!", "The call to unscramble these terms and definitions timed out, so they were not unscrambled!")

        return

def parallelized_blank_fill(question, banned_items, id, domain, context, api_key, model):
    # First retrieve an image embed if it exists
    img_link = get_image_embed_links_if_any(question["Question"])
    
    # convert the banned items into a stringed list
    banned_list = ", ".join(banned_items)
    
    if question["T" + str(id)] == "" and question["D" + str(id)] != "":
        # Get the response from the GPT
        try:
            system_content = SYSTEM_INTRO + REQUEST_EXAMPLE_GENERATE_TERM.replace("{X}", banned_list) + SYSTEM_CONCLUSION + (IMG_SPECIFICATION if img_link != "" else "")
            response_content =  REQUEST_GENERATE_TERM + "Question: "+ remove_backticked_imbeds(question["Question"]) + "Definition: " + question["D" + str(id)]
            response = call_gpt(domain, context, api_key, model, system_content, response_content, img_link)
            
            # idealy, the response should just be as simple as the term to use here, so we can 1-to-1 use it without any further parsing!
            # regardless, at least sanitize the response a little bit.
            question["T" + str(id)] = response.strip().replace('\\\\', '\\')
            
            # add to term log to make sure we dont see this term again as a duplicate!
            banned_items += [question["T" + str(id)]]
        except TimeoutError:
            # fail
            messagebox.showerror("GPT Call Fail!", "The call to ChatGPT to fill in the term timed out!")
            return
    if question["D" + str(id)] == "" and question["T" + str(id)] != "":
        # Get the response from the GPT
        try:
            system_content = SYSTEM_INTRO + REQUEST_EXAMPLE_GENERATE_DEFINITION.replace("{X}", banned_list) + SYSTEM_CONCLUSION + (IMG_SPECIFICATION if img_link != "" else "")
            response_content = REQUEST_GENERATE_DEFINITION + "Question: "+ remove_backticked_imbeds(question["Question"]) + "Term: " + question["T" + str(id)]
            response = call_gpt(domain, context, api_key, model, system_content, response_content, img_link)
            
            # idealy, the response should just be as simple as the term to use here, so we can 1-to-1 use it without any further parsing!
            question["D" + str(id)] = response.strip().replace('\\\\', '\\')
            
            # add to term log to make sure we dont see this term again as a duplicate!
            banned_items += [question["D" + str(id)]]
        except TimeoutError:
            # fail
            messagebox.showerror("GPT Call Fail!", "The call to ChatGPT to fill in the definition timed out!")
            return

def fill_matching_options(question, scrambled, domain, context, api_key, model):
    # iterate through the term definition pairs of the question
    matchings = sum(1 for key in question if key.startswith("D"))
    terms = []
    definitions = []
    for i in range(matchings):
        if question["T" + str(i + 1)] != "":
            terms += [question["T" + str(i + 1)]]
        elif question["D" + str(i + 1)] != "":
            definitions += [question["D" + str(i + 1)]]
            
    # if we have no content, generate some terms relating to the question content, then populate those terms with definitions later on.
    if matchings == 0:
        # First retrieve an image embed if it exists
        img_link = get_image_embed_links_if_any(question["Question"])
        
        # Get the response from the GPT
        try:
            system_content = SYSTEM_INTRO + REQUEST_EXAMPLE_INITAL_TERMS + SYSTEM_CONCLUSION + (IMG_SPECIFICATION if img_link != "" else "")
            response_content = REQUEST_INITAL_TERMS + remove_backticked_imbeds(question["Question"])
            response = call_gpt(domain, context, api_key, model, system_content, response_content, img_link)
            
            # Split the response and use that as our new defacto term list
            terms = [term.strip() for term in response.split(',')]
            
            # Create keys in our question according to the number of terms that we have
            for i in range(len(terms)):
                question["T" + str(i + 1)] = terms[i]
                question["D" + str(i + 1)] = ""
            matchings = len(terms)
        except TimeoutError:
            # fail
            messagebox.showerror("GPT Call Fail!", "The call to ChatGPT to generate terms timed out!")
            return
        
    # fill the blanks in parallel based on what component is missing:
    failsafe = 0
    while any(question[f"T{i}"] == "" for i in range(1, matchings + 1)) or any(question[f"D{i}"] == "" for i in range(1, matchings + 1)) and failsafe < 5:
        threads = []
        t_gen = []
        d_gen = []
        for i in range(matchings):
            if question["T" + str(i + 1)] == "" and question["D" + str(i + 1)] != "":
                t_gen.append(i)
                threads.append(threading.Thread(target=parallelized_blank_fill, args=(question, terms, i+1, domain, context, api_key, model)))
                threads[-1].start()
            elif question["D" + str(i + 1)] == "" and question["T" + str(i + 1)] != "":
                d_gen.append(i)
                threads.append(threading.Thread(target=parallelized_blank_fill, args=(question, definitions, i+1, domain, context, api_key, model)))
                threads[-1].start()
        for thread in threads:
            thread.join()
        threads.clear()
        
        # check for duplicates and "remove" them if necessary
        if len(terms) != len(set(terms)):
            dupe_terms = [item for item in set(terms) if terms.count(item) > 1 for _ in range(terms.count(item) - 1)]
            print("Duplicate terms! " + str(dupe_terms))
            for i in range(matchings):
                if question["T" + str(i + 1)] in dupe_terms and i + 1 in t_gen:
                    dupe_terms.remove(question["T" + str(i + 1)])
                    question["T" + str(i + 1)] = ""
        if len(definitions) != len(set(definitions)):
            dupe_terms = [item for item in set(definitions) if definitions.count(item) > 1 for _ in range(definitions.count(item) - 1)] 
            print("Duplicate definitions! " + str(dupe_terms))
            for i in range(matchings):
                if question["D" + str(i + 1)] in dupe_terms and i + 1 in d_gen:
                    dupe_terms.remove(question["D" + str(i + 1)])
                    question["D" + str(i + 1)] = ""
                    
        t_gen.clear()
        d_gen.clear()
        failsafe += 1
        
    if failsafe == 5:
        print("FAILSAFE TRIGGERED! TOO MANY LOOPS!")
        
    # once the terms and definitions are populated, if prompeted unscramble them to be the correct matchings
    if(scrambled):
        unscramble_matching(question, domain, context, api_key, model)


def fill_essay_guidelines(question, domain, context, api_key, model):
    # First retrieve an image embed if it exists
    img_link = get_image_embed_links_if_any(question["Question"])
    
    # determine how the student response will be graded 
    match question['Format']:
        case "0": # Explaination
            format = "Graded as an explaination; "
        case "1": # Code
            format = "Graded as code written in " + question['Language'] + "; "
        case "2": # Proof
            format = "Graded as a proof; "
    
    # Check if we are missing the guidelines to the question
    if question["Guidelines"] == "":
        # Get the response from the GPT
        try:
            system_content = SYSTEM_INTRO + REQUEST_EXAMPLE_GENERATE_ESSAY_GUIDELINES + SYSTEM_CONCLUSION + (IMG_SPECIFICATION if img_link != "" else "")
            response_content = REQUEST_GENERATE_ESSAY_GUIDELINES + format + "Question: "+ remove_backticked_imbeds(question["Question"])
            response = call_gpt(domain, context, api_key, model, system_content, response_content, img_link)
            
            # idealy, the response should just be as simple as the guidelines for which we follow, so we just sanitize and use it directly
            question["Guidelines"] = response.strip().replace('\\\\', '\\')
        except TimeoutError:
            # fail
            messagebox.showerror("GPT Call Fail!", "The call to ChatGPT to generate essay guidelines timed out!")
            return

def fill_multiple_choice_options(question, domain, context, api_key, model):
    # First retrieve an image embed if it exists
    img_link = get_image_embed_links_if_any(question["Question"])

    # Determine what parts of the question we are missing
    # Generate the correct answer(s) and incorrect answer(s) if both are (effectively) missing
    if not "C1" in question.keys() and (not "A1" in question.keys() or ("A1" in question.keys() and not "A2" in question.keys())) :
        # Get the response from the GPT
        try:
            system_content = SYSTEM_INTRO + REQUEST_EXAMPLE_ALL_CHOICES + SYSTEM_CONCLUSION + (IMG_SPECIFICATION if img_link != "" else "")
            response_content = REQUEST_ALL_CHOICES + remove_backticked_imbeds(question["Question"])
            response = call_gpt(domain, context, api_key, model, system_content, response_content, img_link)
            
            # Parse to get the correct and incorrect answers generated
            try:
                # iterate through the output
                lines = response.split('\n')
                cur_key = ""
                num_cor = 0
                num_incor = 0 if not "A1" in question.keys() else 1
                for line in lines:
                    # check if this line contains a correct answer
                    if line.strip().startswith('C~'):
                        # get the new key
                        num_cor += 1
                        cur_key = "C" + str(num_cor)
                        
                        # set this line into the question dictionary
                        question[cur_key] = line[2:].strip()
                    # check if this line contains a incorrect answer
                    elif line.strip().startswith('I~'):
                        # get the new key
                        num_incor += 1
                        cur_key = "A" + str(num_incor)
                        
                        # set this line into the question dictionary
                        question[cur_key] = line[2:].strip()
                    # otherwise is continuation of previous line
                    else:
                        # if we enter here without a key, something went wrong
                        if cur_key == "":
                            raise Exception()
                        
                        # append this line to the previous on a newline
                        question[cur_key] = "\n" + line.strip()
            except Exception as e:
                # prompt error and prevent submit
                messagebox.showerror("GPT Failed!", "The call to ChatGPT to generate the multiple choice options failed: " + str(e))
                print("Failed call to GPT!")
            
                return
        except TimeoutError:
            # fail
            messagebox.showerror("GPT Call Fail!", "The call to ChatGPT to generate the multiple choice options timed out!")
            return
    
    # generate several incorrect answers if they are missing for several correct answers
    elif "C2" in question.keys() and not "A1" in question.keys():
        # get the required content for the prompt
        correct_answers_c = "Correct Answers: "
        correct_count = 0
        for key, value in question.items():
            if "C" in key:
                correct_answers_c += value + ", "
                correct_count += 1
        content = "Question: " + remove_backticked_imbeds(question['Question']) + "\n" + remove_backticked_imbeds(correct_answers_c) + "\n"

        # Get the response from the GPT
        try:
            system_content = SYSTEM_INTRO + REQUEST_EXAMPLE_INCORRECT_CHOICES_MULTIPLE_CORRECT + SYSTEM_CONCLUSION + (IMG_SPECIFICATION if img_link != "" else "")
            response_content = REQUEST_INCORRECT_CHOICES_MULTIPLE_CORRECT.replace("{X}", str(correct_count)) + content
            response = call_gpt(domain, context, api_key, model, system_content, response_content, img_link)
            
            # Parse to get the correct and incorrect answers generated
            try:
                # iterate through the output
                lines = response.split('\n')
                cur_key = ""
                num_cor = 0
                num_incor = 0
                for line in lines:
                    # check if this line contains a incorrect answer
                    if line.strip().startswith('I~'):
                        # get the new key
                        num_incor += 1
                        cur_key = "A" + str(num_incor)
                        
                        # set this line into the question dictionary
                        question[cur_key] = line[2:].strip()
                    # otherwise is continuation of previous line
                    else:
                        # if we enter here without a key, something went wrong
                        if cur_key == "":
                            raise Exception()
                        
                        # append this line to the previous on a newline
                        question[cur_key] = "\n" + line.strip()
            except Exception as e:
                # prompt error and prevent submit
                messagebox.showerror("GPT Failed!", "The call to ChatGPT to generate the multiple choice options failed: " + str(e))
                print("Failed call to GPT!")
            
                return
        except TimeoutError:
            # fail
            messagebox.showerror("GPT Call Fail!", "The call to ChatGPT to generate the multiple choice options timed out!")
            return
        
    # select a correct answer from the incorrect answers if enough incorrect answers are given.
    elif not "C1" in question.keys() and "A2" in question.keys():
        # First, generate the additional content
        answer_choices_c = "Answer Choices:\n"
        choice_num = 1
        for key, value in question.items():
            if "A" in key:
                answer_choices_c += "A" + str(choice_num) + "~ " + value + "\n"
                choice_num += 1
        content = "Question: " + remove_backticked_imbeds(question['Question']) + "\n" + answer_choices_c + "\n"

        # Get the response from the GPT
        try:
            system_content = SYSTEM_INTRO + REQUEST_EXAMPLE_SELECT_CORRECT_CHOICE + SYSTEM_CONCLUSION + (IMG_SPECIFICATION if img_link != "" else "")
            response_content = REQUEST_SELECT_CORRECT_CHOICE + content
            response = call_gpt(domain, context, api_key, model, system_content, response_content, img_link)
            
            choices = []

            # parse
            print(response.split('\n'))
            choice_num = 0
            for key, value in question.items():
                if "A" in key:
                    choices.append(value)
                    choice_num += 1
            while choice_num > 0:
                del question["A" + str(choice_num)]
                choice_num -= 1

            # Parse to get the correct and incorrect answers generated
            try:
                # iterate through the output
                lines = response.split(',')
                cur_key = ""
                choice_num = len(lines)
                for i in range(len(lines) - 1, -1, -1):
                    line = lines[i]
                    # validate this choice
                    if line.strip().startswith('A'):
                        # get the new key
                        cur_key = "C" + str(choice_num)
                        choice_num -= 1
                        
                        # set this line into the question dictionary
                        answer = int(line.strip()[1])
                        question[cur_key] = choices.pop(answer - 1)

                    # otherwise is continuation of previous line
                    else:
                        # if we enter here, we have a problem
                        raise Exception()
            
            except Exception as e:
                # prompt error and prevent submit
                messagebox.showerror("GPT Failed!", "The call to ChatGPT to generate the multiple choice options failed: " + str(e))
                print("Failed call to GPT!")
            
                return
            
            # push the remainder of the list as incorrect answer choices
            choice_num = 1
            for item in choices:
                question["A" + str(choice_num)] = item
                choice_num += 1
        except TimeoutError:
            # fail
            messagebox.showerror("GPT Call Fail!", "The call to ChatGPT to generate the multiple choice options timed out!")
            return
        
    # generate incorrect answers if they are missing for a single correct answer
    elif "C1" in question.keys() and not "A1" in question.keys():
        # first, generate additional content
        content = "Question: " + remove_backticked_imbeds(question["Question"]) + "\nCorrect Answer: " + remove_backticked_imbeds(question["C1"]) + "\n"

        # Get the response from the GPT
        try:
            system_content = SYSTEM_INTRO + REQUEST_EXAMPLE_INCORRECT_CHOICES_SINGLE_CORRECT + SYSTEM_CONCLUSION + (IMG_SPECIFICATION if img_link != "" else "")
            response_content = REQUEST_INCORRECT_CHOICES_SINGLE_CORRECT + content
            response = call_gpt(domain, context, api_key, model, system_content, response_content, img_link)

            # Parse to get the correct and incorrect answers generated
            try:
                # iterate through the output
                lines = response.split('\n')
                cur_key = ""
                num_cor = 0
                num_incor = 0
                for line in lines:
                    # check if this line contains a incorrect answer
                    if line.strip().startswith('I~'):
                        # get the new key
                        num_incor += 1
                        cur_key = "A" + str(num_incor)
                        
                        # set this line into the question dictionary
                        question[cur_key] = line[2:].strip()
                    # otherwise is continuation of previous line
                    else:
                        # if we enter here without a key, something went wrong
                        if cur_key == "":
                            raise Exception()
                        
                        # append this line to the previous on a newline
                        question[cur_key] = "\n" + line.strip()
            
            except Exception as e:
                # prompt error and prevent submit
                messagebox.showerror("GPT Failed!", "The call to ChatGPT to generate the multiple choice options failed: " + str(e))
                print("Failed call to GPT!")
            
                return
        except TimeoutError:
            # fail
            messagebox.showerror("GPT Call Fail!", "The call to ChatGPT to generate the multiple choice options timed out!")
            return

def get_question_sample(bank):
    # get the proportionate amounts of each question type in the bank, while also currating lists for each of the 3 question types
    mc_questions = []
    td_questions = []
    ess_questions = []
    for question in bank:
        match question["Type"]:
            case "MC":
                mc_questions += [question]
            case "TD":
                td_questions += [question]
            case "Ess":
                ess_questions += [question]
    proportions = (len(mc_questions) / len(bank), len(td_questions) / len(bank), len(ess_questions) / len(bank))
    
    # using those proportions, select a type to generate a sample of
    sel_type = ""
    selection = random.random()
    if selection < proportions[0]:
        sel_type = "MC"
    elif selection < proportions[0] + proportions[1]:
        sel_type = "TD"
    else:
        sel_type = "Ess"
    
    # choose the appropriate number of questions to create the sample
    match sel_type:
        case "MC":
            sample = random.sample(mc_questions, min(3, len(mc_questions)))
        case "TD":
            sample = random.sample(td_questions, 1)
        case "Ess":
            sample = random.sample(ess_questions, 1)
    
    # return the sample
    return sample

def batch_generate_questions(bank, count, domain, context, api_key, model):
    # store ALL AI generated questions in an accumulated list
    ai_generated_questions = []
    
    # create as many threads as needed to process the count, using constants on max threads and batch size to determine our thread distribution
    threads = []
    results = []
    while count > 0:
        # get a quota for this iteration
        quota = MAX_THREADS * THREAD_SIZE if count >= MAX_THREADS * THREAD_SIZE else count
        num_of_threads = quota // THREAD_SIZE + 1 if quota < MAX_THREADS * THREAD_SIZE else MAX_THREADS
        size_of_last_thread = quota % THREAD_SIZE if quota < MAX_THREADS * THREAD_SIZE else THREAD_SIZE
        print(count)
        for i in range(num_of_threads):
            # get a sample of either 3 multiple choice questions, 1 TD question, or 1 essay question, where the type of question we sample is determined by the proportion of that question type in the bank
            sample = get_question_sample(bank)
            
            threads.append(threading.Thread(target=generate_questions, args=(sample, THREAD_SIZE if i < num_of_threads - 1 else size_of_last_thread, domain, context, api_key, model, results)))
            threads[i].start()
        for i in range(num_of_threads):
            threads[i].join()
        ai_generated_questions += results
        threads.clear()
        results.clear()
        count -= MAX_THREADS * THREAD_SIZE
        
    # once all of the questions are generated, go back and update the q_indicies of every question
    q_count = len(bank)
    for question in ai_generated_questions:
        q_count += 1
        question["Q"] = str(q_count)
        
    # return the AI generated questions collected
    return ai_generated_questions

def generate_questions(sample, count, domain, context, api_key, model, t_results):    
    # break immediately if we are to generate 0 questions
    if count == 0:
        return []
    
    # get a textual representation of the prompt count to emphasize the amount requested
    if count == 1:
        num_of_questions = "ONE"
    elif count == 2:
        num_of_questions = "TWO"
    elif count == 3:
        num_of_questions = "THREE"
    elif count == 4:
        num_of_questions = "FOUR"
    elif count == 5:
        num_of_questions = "FIVE"
        
    # format the questions of the sample cleanly for the AI to use as an example
    questions = ""
    for question in sample:
        for key, value in question.items():
            # see if we are writing the question
            if key == "Question":
                questions += "Q~ " + remove_backticked_imbeds(value) + "\n"
        # append seperator
        questions += "\n"
        
    # Get the response from the GPT
    try:
        rng_context = get_random_context_segment(context, CONTEXT_CUTOFF)
        system_content = SYSTEM_INTRO + REQUEST_EXAMPLE_QUESTIONS.replace('{X}', str(len(sample))).replace('{Y}', num_of_questions) + SYSTEM_CONCLUSION
        response_content = questions + REQUEST_QUESTIONS.replace('{X}', num_of_questions)
        response = call_gpt(domain, rng_context, api_key, model, system_content, response_content, "")
        
        # using the question component, construct a dummy question object based on what the sample type was
        new_question = {
            "Question": response.split('~')[1],
            "AI Generated": "True",
            "Type": sample[0]["Type"],
            "Explaination": ""
        }
        
        # if this is an essay type question, we inherit settings from the sample taken
        if new_question["Type"] == "Ess":
            new_question["Guidelines"] = ""
            new_question["Format"] = sample[0]["Format"]
            new_question["Language"] = ""
            new_question["Difficulty"] = "1"
        # otherwise, force the question to always prompt as a MC
        else:
            new_question["Forced"] = "1"
        
        # with the bare minimum information acquired, call the respectiving missing content function to fill in the remaining parameters
        match new_question['Type']:
            # to ensure consistency, the context segment that we used to generate this question will be passed along  
            case "MC":
                fill_multiple_choice_options(new_question, domain, rng_context, api_key, model)
                generate_explaination_for_question(new_question, domain, rng_context, api_key, model)
                t_results += [new_question]
            case "TD":
                fill_matching_options(new_question, False, domain, rng_context, api_key, model)
                generate_explaination_for_question(new_question, domain, rng_context, api_key, model)
                t_results += [new_question]
            case "Ess":
                fill_essay_guidelines(new_question, domain, rng_context, api_key, model)
                t_results += [new_question]
    
        return
    except TimeoutError:
        # fail
        messagebox.showerror("GPT Call Fail!", "The call to ChatGPT to generate the multiple choice options timed out!")
        return
    except Exception as e:
        # fail
        print(f"Error: {str(e)}")
        messagebox.showerror("GPT Gen Fail!", "The call to ChatGPT generated an exception, this is expected, but you will likely have less AI questions than anticipated!")
        return

def get_random_numaric_str(unique_arr, character_set, stri, unique):
    # copy the passed string to modify
    num_str = stri
    
    # iterate over the characters of the provided numaric string
    for i in range(len(num_str)):
        # if the character is not in the character set, replace it with a character from the character set
        if not num_str[i] in character_set:
            num_str = num_str[:i] + random.choice(character_set[:-1]) + (num_str[i+1:] if i < len(num_str) - 1 else "")
    
    # perform 100 attempts at generating a unique bit string if specified
    attempts = 0
    while unique and attempts < 100 and num_str in unique_arr:
        # reset number string
        num_str = stri
        
        # iterate over the characters of the provided numaric string
        for i in range(len(num_str)):
            # if the character is not in the character set, replace it with a character from the character set
            if not num_str[i] in character_set:
                num_str = num_str[:i] + random.choice(character_set) + (num_str[i+1:] if i < len(num_str) - 1 else "")
                
    # add the bit string to the unique list
    unique_arr += [num_str]
    
    # return the result
    return num_str
    
def get_random_number(unique_arr, base, min, max, n, unique):
    # validate bounds check
    if n < 1:
        n = 1
        
    # iterate to make the list of numbers
    result = ""
    for i in range(n):
        # attempt to generate the random number
        rand_num = random.randint(min, max)
        
        # perform 100 attempts at generating a unique number if specified
        attempts = 0
        while unique and attempts < 100 and rand_num in unique_arr:
            # attempt to generate the random number
            rand_num = random.randint(max, min)
            attempts += 1
            
        # add the number to the unique list
        num = int_to_base(rand_num, base)
        unique_arr += [num]
        
        # add to result w/ a comma if necessary
        result += num
        if i < n - 1:
            result += ", "
            
    # return the result
    return result

def int_to_base(number, base):
    # correct base
    base = max(min(base, 64), 0)
        
    # define character set
    if base < 36:
        digits = "0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZ"
    else:
        digits = "ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789+/"
    
    # default 0 case
    if number == 0:
        return '0'
    
    result = ''
    negative = False
    
    # remeber negavity
    if number < 0:
        number = -number
        negative = True
    
    # successively divide by the base and encode the remainder as its respective digit according to the digit character set.
    while number > 0:
        number, remainder = divmod(number, base)
        result = digits[remainder] + result
    
    # add the negative back if needed
    if negative:
        result = '-' + result
    
    # return the result
    return result

def seed_permutations_in_question(question_text, domain, context, api_key, model):
    # First retrieve an image embed if it exists
    img_link = get_image_embed_links_if_any(question_text)
    
    # Split the question content into chunks to search for the first img embed
    embeds = question_text.split('```')
    
    # For numbers to track uniqueness, we use a list
    unique_vals = []
    
    # Keep references for 2 seperate lists:
    segments = [""] # Segments of the completed final string
    segment_c = 0
    blank_defaults = [] # The original values of blanks, as well as the "type" of permutateable part they are
    blanks = [] # The blanks to slowly use to guide GPT to constructing a permutated question
    
    for i, chunk in enumerate(embeds):
        if i % 2 != 0:
            # determine what imbed we are using
            header = chunk.split(':')[0]
            if header == "prt":
                # this is a permutatable text part, get the original text and add it to the cut segment
                blank_defaults.append((chunk[4:], "reg"))
                blanks.append("_" * len(chunk[4:]))
                segment_c += 1 # add a new empty segment to catch the static text of the next segement after this blank
                segments += [""] 
            elif header == "mathprt":
                # this is a MATH permutatable part, signify this in the tuple, otherwise its the same as a prt segment
                blank_defaults.append((chunk[8:], "math"))
                blanks.append("_" * len(chunk[8:]))
                segment_c += 1 # add a new empty segment to catch the static text of the next segement after this blank
                segments += [""]
            elif header == "rng":
                # this is a random integer permutable part, which is foramtted as: ```rng:base:min:max:n:unique```
                # get those args and pass them into the randomization function
                try:
                    default_args = ["per_name", 10, 1, 10, 1, False]
                    args = chunk.split(':')
                    argc = len(args)
                    for i in range(1, 6):
                        if i < argc:
                            args[i] = int(args[i]) if i != 5 else ("t" in args[i].strip().lower())
                            continue
                        args += [default_args[i]]
                    args.pop(0)
                    print(args)
                    # automatically fill the blank with the random number
                    segments[segment_c] += get_random_number(unique_vals, *args)
                except Exception as e:
                    print("Malformated Embed " + str(e))
                    segments[segment_c] += str(0)
            elif header == "rnstr":
                # this is a binary permutable part, which is foramtted as: ```rnstr:base:str:unique```
                # get those args and pass them into the randomization function
                try:
                    default_args = ["per_name", 10, "-", False]
                    args = chunk.split(':')
                    argc = len(args)
                    for i in range(1, 4):
                        if i < argc:
                            if i == 1:
                                args[i] = int(args[i])
                            elif i == 3:
                                args[i] = ("t" in args[i].strip().lower())
                            continue
                        args += [default_args[i]]
                    args.pop(0)
                    
                    base = args.pop(0)
                    # correct base
                    base = min(max(base, 0), 64)
                        
                    # define character set
                    if base < 36:
                        digits = "0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZ"
                    else:
                        digits = "ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789+/"

                    # automatically fill the blank with the random number
                    segments[segment_c] += get_random_numaric_str(unique_vals, digits[0:base] + "_", *args)
                except Exception as e:
                    print("Malformated Embed " + str(e))
                    segments[segment_c] += str(0)
            else:
                # readd embed markings if not an image embed
                if not chunk.startswith("img:"):
                    segments[segment_c] += "```" + chunk + "```"
        else:
            # log segment as is
            segments[segment_c] += chunk
    
    # iterate getting responses from our client to fill in the blanks as they appear.
    # blanks already filled are considered a singular segment, the current blank appears blank and is requested to be filled, where all remaining blanks are filled with their default values.
    while not len(blanks) == 0:
        # sew a question together consisting of a blank where it first sees fit (between the first and second segment), followed by the remaining segments joined using the original default values as they appear
        question_c = ""
        comb_segments = (segments.pop(0), segments.pop(0))
        cur_blank = blank_defaults.pop(0)
        question_c = comb_segments[0] + blanks.pop(0) + comb_segments[1]
        for i, def_blank in enumerate(blank_defaults):
            if i < len(segments):
                question_c += def_blank[0] + segments[i]
            else:
                question_c += def_blank[0]
        
        # Get the response from the GPT
        try:
            system_content = SYSTEM_INTRO + (REQUEST_EXAMPLE_FILL_IN_BLANKS if cur_blank[1] == "reg" else REQUEST_EXAMPLE_FILL_IN_MATH) + SYSTEM_CONCLUSION + (IMG_SPECIFICATION if img_link != "" else "")
            response_content = REQUEST_FILL_IN_BLANKS + question_c
            response = call_gpt(domain, context, api_key, model, system_content, response_content, img_link)
            
            try:
                # pick one of the 5 options at basically random
                options = response.split('\n')
                print(options)
                best_option = random.choice(options)

                if cur_blank[1] == "reg":
                    segments.insert(0, comb_segments[0] + best_option.strip() + comb_segments[1])
                else:
                    segments.insert(0, comb_segments[0] + "```math:" + best_option.strip() + "```" + comb_segments[1])
            except Exception as e:
                # prompt error and prevent submit
                messagebox.showerror("GPT Failed!", "The call to ChatGPT to generate the multiple choice options failed: " + str(e))
                print("Failed call to GPT!")
                
                return None
        except TimeoutError:
            # fail
            messagebox.showerror("GPT Call Fail!", "The call to ChatGPT to generate the multiple choice options timed out!")
            return
        
    # return the results
    return segments[0]
    
def generate_explaination_for_question(question, domain, context, api_key, model):
    # First retrieve an image embed if it exists
    img_link = get_image_embed_links_if_any(question["Question"])

    # switch based on the question content
    match question["Type"]:
        case "MC":
            # get the required content for the prompt
            correct_answers_c = "Correct Answers: "
            for key, value in question.items():
                if "C" in key:
                    correct_answers_c += value + ", "
            content = "Question: " + remove_backticked_imbeds(question['Question']) + "\n" + remove_backticked_imbeds(correct_answers_c) + "\n"

            # Get the response from the GPT
            try:
                system_content = SYSTEM_INTRO + REQUEST_EXAMPLE_EXPLAINATIONS_MC + SYSTEM_CONCLUSION + (IMG_SPECIFICATION if img_link != "" else "")
                response_content = REQUEST_EXPLAINATIONS_MC + content
                response = call_gpt(domain, context, api_key, model, system_content, response_content, img_link)
                
            except TimeoutError:
                # fail
                messagebox.showerror("GPT Call Fail!", "The call to ChatGPT to generate the multiple choice options timed out!")
                return

        case "TD":
            # get the required content for the prompt
            matchings = "Matchings: "
            item_c = sum(1 for key in question if key.startswith("D"))
            for i in range(item_c):
                matchings += question["T" + str(i + 1)] + "-   " + question["D" + str(i + 1)] + "\t"
            content = "Question: " + remove_backticked_imbeds(question['Question']) + "\n" + matchings + "\n"

            # Get the response from the GPT
            try:
                system_content = SYSTEM_INTRO + REQUEST_EXAMPLE_EXPLAINATIONS_TD + SYSTEM_CONCLUSION + (IMG_SPECIFICATION if img_link != "" else "")
                response_content = REQUEST_EXPLAINATIONS_TD + content
                response = call_gpt(domain, context, api_key, model, system_content, response_content, img_link)
                
            except TimeoutError:
                # fail
                messagebox.showerror("GPT Call Fail!", "The call to ChatGPT to generate the multiple choice options timed out!")
                return
        case "Ess":
            # explainations for essay questions are generated at grade time, so we don't need to generate any explaination
            question['Explaination'] = "Explainations for essay question to be generated at grade time, so none will be provided :)"
            
            return

    # wrap in try/except to flag an error if the api fails
    try:
        # push to explaination, sanitize
        if "E~" in response:
            question['Explaination'] = response[2:].strip().replace('\\\\', '\\').replace("\n", " ")
        else:
            question['Explaination'] = response.strip().replace('\\\\', '\\').replace("\n", " ")
        print(question['Explaination'])
    except Exception as e:
        # prompt error and prevent submit
        messagebox.showerror("GPT Failed!", "The call to ChatGPT to generate the multiple choice options failed: " + str(e))
        print("Failed call to GPT!")
        return

def test_api_key(api_key):
    # test api key
    try:
        result = call_gpt("N/A", "N/A", api_key, 0, "Only respond to the user with 'OK'", "This is a test prompt, please respond with ONLY 'OK'", "")
    except Exception as e:
        return False
    
    if not "OK" in result:
        return False
    return True

def import_from_quizzy_file():
    # get the path to the file we want to open
    print("Requesting Quizzy file location...")
    bank_path = filedialog.askopenfilename(defaultextension='.qizy', filetypes=[("Quizzy Files", "*.qizy")])
    
    # check if the path is valid
    if bank_path:
        try:
            # open the file for parsing
            with open(bank_path, 'r', encoding='utf-8') as import_file:
                json_data = import_file.read()
                
            # convert from JSON to dictionary
            py_obj = json.loads(json_data)
            extracted_values = py_obj
            
            # check if all contents are there
            domain = extracted_values['Domain']
            print("Read domain successfully!")
            data = extracted_values['Data']
            print("Read data successfully!")
            
            # decode the context segement, as it is encrypted.
            cipher_suite = Fernet(CONTEXT_CRYPTO_KEY)
            decoded_context = cipher_suite.decrypt(extracted_values['Context'].encode("utf-8")).decode('utf-8')
            context = decoded_context
            print("Read context successfully!")
            print("All contents found SUCCESSFULLY!")
            
            # iterate through the question objects and modify/validate them to be correct
            valid = all(validate_question_object(obj) for obj in data)
            if not valid:
                print("Failed to read Quizzy file. Is it formatted correctly?")
            else:
                # return the domain, context, and data segment as a tuple
                print("Data segment VALIDATED!")
                print("Successfully read Quizzy file!")
                return (domain, context, data, bank_path)
        except Exception as e:
            print("Failed to read Quizzy file. Is it formatted correctly?")
    else:
        print("Failed to find Quizzy file. Is the path correct?")
        
    return ("No domain provided", "No context provided", [], "")

def export_to_quizzy_file(domain, context, data):
    # get the path to the file we want to save
    print("Requesting location to store Quizzy file at...")
    save_path = filedialog.asksaveasfile(initialfile=f'{domain}.qizy', mode='w', defaultextension='.qizy', filetypes=[("Quizzy Files", "*.qizy"), ("All Files", "*.*")]).name # File dialog to get a file path to save at

    # fail if the file path is invalid
    if not save_path:
        print("Failed to find path. Was it written correctly?")
        messagebox.showerror("Error", "There was an issue saving this deck to that location!")
        return False
    
    # encode the context segement so that its information is compacted and obscurred when stored as a file to still allow quizzy files to be edited at the file level.
    cipher_suite = Fernet(CONTEXT_CRYPTO_KEY)
    encoded_context = cipher_suite.encrypt(context.encode("utf-8")).decode('utf-8')
    
    # create a save data object
    save = {"Domain":domain,
            "Context":encoded_context,
            "Data":data}
    
    # export to JSON file
    json_data = json.dumps(save, indent=4)
    with open(save_path, 'w', encoding='utf-8') as export:
        # write to destination
        export.write(json_data)
        
    print("Successfully saved Quizzy file")
    return True

def get_text_context_from_file(files):
    # use a variable to collate all text content
    text_content = ""
    
    # process the text content of the file differently depending on what type it is
    for file_path in files:
        extension = file_path.split(".")[-1]
        match extension:
            case "txt":
                text_content += open(file_path, "r").read()
            case "pdf":
                doc = fitz.open(file_path)
                for page in doc:
                    text_content += page.get_text() + "\n\n"
                doc.close()
            case "docx":
                text_content += docx2txt.process(file_path)
            case "pptx":
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content += shape.text + "\n"
            case _:
                text_content += open(file_path, "r").read()
            
    return text_content
